#!/usr/bin/env python3
"""
Generate PowerPoint presentation for TankAlarm 112025 Opta Setup
This script creates a comprehensive setup guide with screenshots and step-by-step instructions.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Presentation Configuration
PRESENTATION_TITLE = "TankAlarm 112025 - Opta Setup Guide"
SUBTITLE = "Hardware Wiring & Software Installation"
OUTPUT_FILE = "TankAlarm_112025_Setup_Guide.pptx"

# Slide Dimensions (in inches)
SLIDE_WIDTH = 10.0
SLIDE_HEIGHT = 7.5

# Color Scheme (Arduino-themed)
COLOR_ARDUINO_BLUE = RGBColor(0, 32, 96)      # Dark blue for titles
COLOR_SECTION_BLUE = RGBColor(0, 120, 215)   # Bright blue for sections
COLOR_BLACK = RGBColor(0, 0, 0)              # Black for body text

# Screenshot paths (relative to script location)
SCREENSHOT_BASE = Path(__file__).parent

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style the title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_ARDUINO_BLUE

def add_section_slide(prs, section_title):
    """Add section divider slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = section_title
    
    # Style the section title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_SECTION_BLUE

def add_content_slide(prs, title, content_items, image_path=None):
    """Add content slide with bullet points and optional image"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Style title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # Add content
    if image_path and os.path.exists(image_path):
        # Content on left, image on right
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(4.5)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(16)
            p.space_before = Pt(6)
        
        # Add image on right
        img_left = Inches(5.5)
        img_top = Inches(1.5)
        img_width = Inches(4)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    else:
        # Full width content
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(8)

def add_image_slide(prs, title, image_path, caption=None):
    """Add slide with large image and optional caption"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add image if it exists
    if os.path.exists(image_path):
        img_top = Inches(1.2)
        img_left = Inches(1)
        img_width = Inches(8)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    else:
        # Add placeholder text if image doesn't exist
        placeholder_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        tf = placeholder_box.text_frame
        tf.text = f"Screenshot placeholder:\n{os.path.basename(image_path)}"
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
    
    # Add caption if provided
    if caption:
        caption_top = Inches(6.8)
        caption_box = slide.shapes.add_textbox(Inches(1), caption_top, Inches(8), Inches(0.6))
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    
    # Title slide
    add_title_slide(prs, PRESENTATION_TITLE, SUBTITLE)
    
    # Section 1: Overview
    add_section_slide(prs, "1. System Overview")
    
    add_content_slide(prs, "What is TankAlarm 112025?", [
        "Cellular tank level monitoring system",
        "Arduino Opta + Blues Wireless cellular connectivity",
        "Client devices monitor tank levels",
        "Server aggregates data and provides web dashboard",
        "SMS and email alerts for alarms",
        "No SD card required - uses internal flash storage",
        "Remote configuration via web interface"
    ])
    
    add_content_slide(prs, "System Components", [
        "Client Devices (Tank Monitors):",
        "  • Arduino Opta Lite",
        "  • Blues Wireless for Opta",
        "  • Arduino Opta Ext A0602 (analog expansion)",
        "  • Level sensors (analog/4-20mA)",
        "",
        "Server Device (Data Aggregation):",
        "  • Arduino Opta Lite",
        "  • Blues Wireless for Opta",
        "  • Ethernet connection",
        "  • Web dashboard for monitoring & configuration"
    ])
    
    # Section 2: Hardware Requirements
    add_section_slide(prs, "2. Hardware Requirements")
    
    add_content_slide(prs, "Client Hardware", [
        "Required Components:",
        "• Arduino Opta Lite ($95)",
        "• Blues Wireless for Opta ($49)",
        "• Arduino Pro Opta Ext A0602 ($56)",
        "  (Analog/4-20mA expansion module)",
        "• USB-C cable for programming",
        "• 24V DC power supply",
        "• Tank level sensors (analog/digital/4-20mA)",
        "",
        "Purchase Links:",
        "• store-usa.arduino.cc/products/opta-lite",
        "• shop.blues.com"
    ])
    
    add_content_slide(prs, "Server Hardware", [
        "Required Components:",
        "• Arduino Opta Lite ($95)",
        "• Blues Wireless for Opta ($49)",
        "• USB-C cable for programming",
        "• Ethernet cable (RJ45)",
        "• 12-24V DC power supply",
        "",
        "Network Requirements:",
        "• Access to local network (intranet)",
        "• DHCP enabled (or configure static IP)",
        "• Port 80 accessible for web dashboard"
    ])
    
    # Section 3: Arduino IDE Setup
    add_section_slide(prs, "3. Arduino IDE Installation")
    
    add_content_slide(prs, "Install Arduino IDE", [
        "Download Arduino IDE:",
        "• Visit: arduino.cc/en/software",
        "• Download Arduino IDE 2.0 or later",
        "• Install for your OS (Windows/Mac/Linux)",
        "• Launch Arduino IDE",
        "",
        "First-time setup:",
        "• Accept any firewall prompts",
        "• IDE will download board indexes on first launch",
        "• Wait for initialization to complete"
    ])
    
    add_content_slide(prs, "Install Opta Board Support", [
        "Add Arduino Opta to IDE:",
        "1. Open Tools → Board → Boards Manager",
        "2. Search: 'Arduino Mbed OS Opta Boards'",
        "3. Click Install (may take several minutes)",
        "4. Wait for download and installation",
        "",
        "Verify Installation:",
        "• Tools → Board → Arduino Mbed OS Opta Boards",
        "• You should see 'Arduino Opta' option",
        "• If not visible, restart Arduino IDE"
    ])
    
    # Section 4: Library Installation
    add_section_slide(prs, "4. Required Libraries")
    
    add_content_slide(prs, "Install Required Libraries", [
        "Open Library Manager:",
        "• Tools → Manage Libraries",
        "",
        "Install These Libraries:",
        "1. ArduinoJson (v7.x or later)",
        "   Search: 'ArduinoJson'",
        "   Install: 'ArduinoJson by Benoit Blanchon'",
        "",
        "2. Blues Wireless Notecard (latest)",
        "   Search: 'Notecard'",
        "   Install: 'Blues Wireless Notecard by Blues Inc.'",
        "",
        "3. Built-in (no install needed):",
        "   • Ethernet, Wire, LittleFS"
    ])
    
    # Section 5: Blues Notecard Setup
    add_section_slide(prs, "5. Blues Notehub Configuration")
    
    add_content_slide(prs, "Create Blues Notehub Account", [
        "Sign up for Blues Notehub:",
        "1. Visit: notehub.io",
        "2. Click 'Sign Up' (free account available)",
        "3. Verify email address",
        "4. Log in to Notehub dashboard",
        "",
        "Create a Product:",
        "1. Click '+ New Product'",
        "2. Name: 'TankAlarm System' (or your choice)",
        "3. Note the Product UID",
        "   Format: com.company.product:project",
        "4. You'll need this for device configuration"
    ])
    
    add_content_slide(prs, "Create Device Fleets", [
        "Fleet-based communication setup:",
        "",
        "Create Server Fleet:",
        "1. Go to Fleets → Create Fleet",
        "2. Name: 'tankalarm-server'",
        "3. Description: 'Base station server'",
        "",
        "Create Client Fleet:",
        "1. Go to Fleets → Create Fleet",
        "2. Name: 'tankalarm-clients'",
        "3. Description: 'Field monitoring devices'",
        "",
        "Benefits:",
        "• No manual route configuration needed",
        "• Automatic device-to-device messaging",
        "• Easy to scale with more devices"
    ])
    
    add_content_slide(prs, "Provision Notecards", [
        "Add Notecards to your Product:",
        "",
        "For each device:",
        "1. Power on the Arduino Opta with Notecard",
        "2. In Notehub, go to Devices",
        "3. The Notecard should appear automatically",
        "4. Click 'Claim Device'",
        "5. Assign to appropriate fleet:",
        "   • Server → 'tankalarm-server' fleet",
        "   • Clients → 'tankalarm-clients' fleet",
        "",
        "Note the Device UID:",
        "• Format: dev:xxxxxxxxxxxxxxxxxx",
        "• You'll need this for configuration"
    ])
    
    # Section 6: Client Wiring
    add_section_slide(prs, "6. Client Device Wiring")
    
    add_content_slide(prs, "Client Wiring Overview", [
        "Connection Steps:",
        "1. Install Blues Wireless for Opta on Opta Lite",
        "2. Connect Opta Ext A0602 expansion module",
        "3. Wire analog sensors to expansion inputs",
        "4. Connect 24V DC power supply",
        "",
        "Sensor Wiring Options:",
        "• Analog voltage sensors (0-10V)",
        "• 4-20mA current loop sensors",
        "• Digital pulse sensors",
        "",
        "Important:",
        "• Do not exceed voltage ratings",
        "• Verify sensor polarity",
        "• Use shielded cable for long sensor runs"
    ])
    
    add_content_slide(prs, "Analog Sensor Connections", [
        "Using Opta Ext A0602 Expansion:",
        "",
        "For each tank sensor:",
        "1. Connect sensor signal to input terminal",
        "   • Channel 0-7 (up to 8 sensors)",
        "2. Connect sensor ground to GND",
        "3. Provide sensor excitation voltage if needed",
        "",
        "Configuration:",
        "• Set sensor type in web interface",
        "• Configure input pin mapping",
        "• Set calibration values",
        "",
        "Supported Sensor Types:",
        "• Analog voltage (ratiometric)",
        "• 4-20mA current loop",
        "• Ultrasonic level sensors"
    ])
    
    # Section 7: Server Wiring
    add_section_slide(prs, "7. Server Device Wiring")
    
    add_content_slide(prs, "Server Wiring Overview", [
        "Connection Steps:",
        "1. Install Blues Wireless for Opta on Opta Lite",
        "2. Connect Ethernet cable to RJ45 port",
        "3. Connect 12-24V DC power supply",
        "4. No sensor expansion needed for server",
        "",
        "Network Connection:",
        "• Connect to local network switch/router",
        "• DHCP will assign IP address",
        "• Or configure static IP in code",
        "",
        "Power Requirements:",
        "• 12-24V DC, 1A minimum",
        "• Use industrial-grade power supply",
        "• Consider UPS for critical installations"
    ])
    
    # Section 8: Firmware Upload - Client
    add_section_slide(prs, "8. Upload Client Firmware")
    
    add_content_slide(prs, "Open Client Sketch", [
        "Load the client code:",
        "1. Download/clone repository from GitHub",
        "2. Navigate to folder:",
        "   TankAlarm-112025-Client-BluesOpta/",
        "3. Open file:",
        "   TankAlarm-112025-Client-BluesOpta.ino",
        "4. Arduino IDE will open the sketch",
        "",
        "Configure Product UID:",
        "1. Find this line near the top:",
        "   #define PRODUCT_UID",
        "2. Replace with your Notehub Product UID",
        "3. Format: \"com.company.product:project\"",
        "4. Save the file (Ctrl+S)"
    ])
    
    add_content_slide(prs, "Compile and Upload Client", [
        "Select Board:",
        "1. Tools → Board → Arduino Mbed OS Opta Boards",
        "2. Select 'Arduino Opta'",
        "3. Tools → Port → Select Opta COM port",
        "",
        "Compile:",
        "1. Click Verify (✓) button",
        "2. Wait for compilation",
        "3. Check for errors in console",
        "",
        "Upload:",
        "1. Click Upload (→) button",
        "2. Wait for upload to complete",
        "3. Success: 'Upload complete' message",
        "",
        "Typical memory usage: 60-80% flash, 40-60% RAM"
    ])
    
    add_content_slide(prs, "Verify Client Operation", [
        "Open Serial Monitor:",
        "1. Tools → Serial Monitor (Ctrl+Shift+M)",
        "2. Set baud rate: 115200",
        "3. Press Reset button on Opta",
        "",
        "Expected Output:",
        "• 'TankAlarm 112025 Client - Blues Opta'",
        "• 'LittleFS initialized'",
        "• 'Notecard UID: dev:xxxxx...'",
        "• 'Configuration loaded from flash'",
        "• Hardware requirements listed",
        "• Initial sensor readings",
        "",
        "If errors occur:",
        "• Check library installation",
        "• Verify board selection",
        "• Check Notecard connection"
    ])
    
    # Section 9: Firmware Upload - Server
    add_section_slide(prs, "9. Upload Server Firmware")
    
    add_content_slide(prs, "Open Server Sketch", [
        "Load the server code:",
        "1. Navigate to folder:",
        "   TankAlarm-112025-Server-BluesOpta/",
        "2. Open file:",
        "   TankAlarm-112025-Server-BluesOpta.ino",
        "3. Arduino IDE will open the sketch",
        "",
        "Configure Product UID:",
        "1. Find this line near the top:",
        "   #define SERVER_PRODUCT_UID",
        "2. Replace with your Notehub Product UID",
        "3. Must match client Product UID",
        "4. Save the file (Ctrl+S)"
    ])
    
    add_content_slide(prs, "Compile and Upload Server", [
        "Follow same steps as client:",
        "1. Select Board: Arduino Opta",
        "2. Select Port: Opta COM port",
        "3. Click Verify (✓)",
        "4. Click Upload (→)",
        "",
        "Note:",
        "• Server sketch is larger than client",
        "• May use 70-85% flash memory",
        "• This is normal and expected",
        "",
        "After upload:",
        "1. Disconnect USB (optional)",
        "2. Connect Ethernet cable",
        "3. Apply power (12-24V DC)"
    ])
    
    add_content_slide(prs, "Verify Server Operation", [
        "Check Serial Monitor (if USB connected):",
        "• 'TankAlarm 112025 Server - Blues Opta'",
        "• 'LittleFS initialized'",
        "• 'Notecard UID: dev:xxxxx...'",
        "• 'Ethernet connected'",
        "• 'IP address: xxx.xxx.xxx.xxx'",
        "• 'Server listening on port 80'",
        "",
        "Note the IP Address!",
        "• You'll need this to access web dashboard",
        "",
        "If no serial access:",
        "• Check your router's DHCP client list",
        "• Look for 'Arduino Opta' device",
        "• Try common IP ranges: 192.168.1.x"
    ])
    
    # Section 10: Website - Dashboard
    add_section_slide(prs, "10. Server Web Dashboard")
    
    # Add screenshot slides for server dashboard
    dashboard_img = str(SCREENSHOT_BASE / "TankAlarm-112025-Server-BluesOpta" / "screenshots" / "dashboard.png")
    add_image_slide(prs, "Server Dashboard - Overview", dashboard_img, 
                   "Real-time tank levels, alarm status, and last update times")
    
    add_content_slide(prs, "Access the Web Dashboard", [
        "Open web browser on your network:",
        "1. Enter server IP address:",
        "   http://192.168.1.xxx/",
        "2. Dashboard should load immediately",
        "",
        "Dashboard Features:",
        "• Real-time tank level display",
        "• Current alarm status",
        "• Last update timestamp for each client",
        "• Client configuration interface",
        "• Server settings management",
        "• Historical data view",
        "",
        "Bookmark this page for easy access!"
    ])
    
    # Section 11: Website - Add Sensor/Alarm
    add_section_slide(prs, "11. Adding Sensors & Alarms")
    
    client_console_img = str(SCREENSHOT_BASE / "TankAlarm-112025-Server-BluesOpta" / "screenshots" / "client-console.png")
    add_image_slide(prs, "Client Configuration Console", client_console_img,
                   "Configure client devices through the web interface")
    
    add_content_slide(prs, "Configure Client Devices", [
        "From the server dashboard:",
        "1. Scroll to 'Update Client Config' section",
        "2. Select client from dropdown menu",
        "   (auto-populated from received telemetry)",
        "",
        "Client Information:",
        "• Device UID shown",
        "• Last contact time displayed",
        "• Current configuration loaded",
        "",
        "If client not in dropdown:",
        "• Wait for first telemetry (30 min default)",
        "• Verify fleet assignment in Notehub",
        "• Check client serial console for errors"
    ])
    
    pin_setup_img = str(SCREENSHOT_BASE / "TankAlarm-112025-Server-BluesOpta" / "screenshots" / "client-console-pin-setup.png")
    add_image_slide(prs, "Sensor Pin Configuration", pin_setup_img,
                   "Set up sensor types and pin assignments for each tank")
    
    add_content_slide(prs, "Add a New Tank Sensor", [
        "In client configuration form:",
        "",
        "1. Basic Settings:",
        "   • Tank ID: A-H (single letter)",
        "   • Tank Name: Descriptive name",
        "   • Tank Number: 1-8",
        "",
        "2. Sensor Configuration:",
        "   • Sensor Type: analog/digital/4-20mA",
        "   • Primary Pin: 0-7 (expansion channel)",
        "   • Secondary Pin: -1 (or redundant sensor)",
        "   • Loop Channel: -1 (or 4-20mA channel)",
        "",
        "3. Tank Parameters:",
        "   • Height: Total tank height in inches",
        "   • High Alarm: Upper threshold in inches",
        "   • Low Alarm: Lower threshold in inches"
    ])
    
    add_content_slide(prs, "Configure Alarm Settings", [
        "For each tank:",
        "",
        "Enable Alarms:",
        "☑ Daily Report - Include in daily summary",
        "☑ Alarm SMS - Send SMS on alarm condition",
        "☑ Upload - Include in telemetry",
        "",
        "Alarm Behavior:",
        "• High alarm triggers when level exceeds threshold",
        "• Low alarm triggers when level drops below threshold",
        "• SMS sent to server's contact list",
        "• Alarm persists until condition clears",
        "",
        "Testing:",
        "• Set test thresholds initially",
        "• Verify SMS delivery",
        "• Adjust to actual operational thresholds"
    ])
    
    add_content_slide(prs, "Send Configuration to Client", [
        "After configuring all settings:",
        "",
        "1. Review configuration:",
        "   • Site name and device label",
        "   • Server fleet: 'tankalarm-server'",
        "   • Sample interval (seconds)",
        "   • All tank configurations",
        "",
        "2. Click 'Send Config to Client' button",
        "",
        "3. Server will:",
        "   • Create configuration note",
        "   • Send via Blues Notehub",
        "   • Target specific client device",
        "",
        "4. Client will:",
        "   • Receive configuration",
        "   • Validate and apply settings",
        "   • Save to internal flash",
        "   • Report back via serial console"
    ])
    
    # Section 12: Website - Other Features
    contacts_img = str(SCREENSHOT_BASE / "TankAlarm-112025-Server-BluesOpta" / "screenshots" / "contacts.png")
    add_image_slide(prs, "SMS Contact Management", contacts_img,
                   "Configure phone numbers for SMS alerts")
    
    calibration_img = str(SCREENSHOT_BASE / "TankAlarm-112025-Server-BluesOpta" / "screenshots" / "calibration.png")
    add_image_slide(prs, "Sensor Calibration", calibration_img,
                   "Calibrate sensors for accurate level readings")
    
    historical_img = str(SCREENSHOT_BASE / "TankAlarm-112025-Server-BluesOpta" / "screenshots" / "historical.png")
    add_image_slide(prs, "Historical Data View", historical_img,
                   "View historical tank levels and trends")
    
    server_settings_img = str(SCREENSHOT_BASE / "TankAlarm-112025-Server-BluesOpta" / "screenshots" / "server-settings.png")
    add_image_slide(prs, "Server Settings", server_settings_img,
                   "Configure server name, fleets, and daily report settings")
    
    # Section 13: Blues Notehub Fleet Management
    add_section_slide(prs, "12. Blues Notehub Device Management")
    
    add_content_slide(prs, "Verify Devices in Notehub", [
        "Check device connectivity:",
        "",
        "1. Log in to notehub.io",
        "2. Select your Product",
        "3. Go to Devices tab",
        "",
        "You should see:",
        "• All Notecards listed",
        "• Status: Connected/Online",
        "• Last activity timestamp",
        "• Fleet assignments",
        "",
        "For each device:",
        "• Click device name to view details",
        "• Check 'Fleets' section",
        "• Verify correct fleet assignment",
        "• Review recent activity in Events"
    ])
    
    add_content_slide(prs, "Monitor Note Traffic", [
        "View communication between devices:",
        "",
        "1. In Notehub, go to Events",
        "2. Filter by device or note type",
        "",
        "Client Notes (outbound to server):",
        "• fleet.tankalarm-server:telemetry.qi",
        "• fleet.tankalarm-server:alarm.qi",
        "• fleet.tankalarm-server:daily.qi",
        "",
        "Server Notes (outbound to clients):",
        "• device:<uid>:config.qi",
        "",
        "Healthy System Shows:",
        "• Regular telemetry from clients",
        "• Config updates when sent",
        "• Recent timestamps (< 1 hour)"
    ])
    
    add_content_slide(prs, "Troubleshooting Connectivity", [
        "If devices not communicating:",
        "",
        "Check Fleet Assignments:",
        "1. Devices → Select device",
        "2. Verify fleet in 'Fleets' section",
        "3. Reassign if incorrect",
        "",
        "Verify Notecard Status:",
        "1. Check for cellular signal",
        "2. Review Events for errors",
        "3. Check SIM card activation",
        "",
        "Force Sync:",
        "1. In Notehub, go to device",
        "2. Click 'Sync Now'",
        "3. Wait for sync to complete",
        "",
        "Check Client Serial Console:",
        "• Should show sync activity",
        "• Look for Notecard errors"
    ])
    
    # Section 14: Testing
    add_section_slide(prs, "13. System Testing")
    
    add_content_slide(prs, "Test Sensor Readings", [
        "Verify sensors are working:",
        "",
        "1. Client Serial Console:",
        "   • Shows real-time sensor readings",
        "   • Values should be reasonable",
        "   • Check units (inches)",
        "",
        "2. Server Dashboard:",
        "   • Wait for telemetry (sample interval)",
        "   • Tank levels should appear",
        "   • Timestamps should update",
        "",
        "3. Calibration:",
        "   • Measure actual tank level",
        "   • Compare to displayed value",
        "   • Adjust calibration if needed",
        "   • Use calibration interface on server"
    ])
    
    add_content_slide(prs, "Test Alarm Functionality", [
        "Verify alarms work correctly:",
        "",
        "1. Set Test Threshold:",
        "   • Configure low alarm above current level",
        "   • Or high alarm below current level",
        "   • Send config to client",
        "",
        "2. Wait for Alarm:",
        "   • Client detects alarm condition",
        "   • Sends alarm note to server",
        "   • Check server dashboard for alarm",
        "",
        "3. Verify SMS (if configured):",
        "   • SMS should be sent",
        "   • Check received message content",
        "   • Verify phone numbers correct",
        "",
        "4. Reset Thresholds:",
        "   • Set to actual operational values",
        "   • Send updated config"
    ])
    
    add_content_slide(prs, "Test Daily Reports", [
        "Verify daily reporting:",
        "",
        "1. Configure Report Time:",
        "   • Set time in server settings",
        "   • Usually early morning (5-7 AM)",
        "",
        "2. For Testing:",
        "   • Set report time to near future",
        "   • Wait for scheduled time",
        "",
        "3. Verify Report:",
        "   • Client sends daily note",
        "   • Server processes report",
        "   • Email sent (if configured)",
        "",
        "4. Production:",
        "   • Set to desired daily time",
        "   • Monitor for first few days",
        "   • Adjust as needed"
    ])
    
    # Section 15: Maintenance
    add_section_slide(prs, "14. Ongoing Maintenance")
    
    add_content_slide(prs, "Regular Monitoring", [
        "Daily tasks:",
        "• Check dashboard for current levels",
        "• Verify all clients reporting",
        "• Review any active alarms",
        "",
        "Weekly tasks:",
        "• Review historical trends",
        "• Check for missed reports",
        "• Verify Blues Notehub connectivity",
        "",
        "Monthly tasks:",
        "• Review Blues data usage",
        "• Check sensor calibration",
        "• Test alarm SMS delivery",
        "• Verify contact phone numbers current",
        "",
        "Best Practices:",
        "• Bookmark server dashboard",
        "• Set up mobile access for alerts",
        "• Document any configuration changes"
    ])
    
    add_content_slide(prs, "Updating Configuration", [
        "To update client settings:",
        "",
        "1. Access server web dashboard",
        "2. Select client from dropdown",
        "3. Modify desired settings:",
        "   • Tank parameters",
        "   • Alarm thresholds",
        "   • Sample interval",
        "   • Server fleet name",
        "4. Click 'Send Config to Client'",
        "5. Verify receipt in client console",
        "",
        "Configuration is saved to:",
        "• Client internal flash (persistent)",
        "• Server internal flash (for reference)",
        "",
        "Changes take effect:",
        "• Immediately after client receives config",
        "• No need to restart or recompile"
    ])
    
    # Final slide
    add_content_slide(prs, "Additional Resources", [
        "Documentation:",
        "• GitHub Repository:",
        "  github.com/SenaxInc/ArduinoSMSTankAlarm",
        "",
        "• Installation Guides:",
        "  - TankAlarm-112025-Client-BluesOpta/INSTALLATION.md",
        "  - TankAlarm-112025-Server-BluesOpta/INSTALLATION.md",
        "",
        "• Fleet Setup Guide:",
        "  - TankAlarm-112025-Server-BluesOpta/FLEET_SETUP.md",
        "",
        "External Resources:",
        "• Arduino Opta: docs.arduino.cc/hardware/opta",
        "• Blues Wireless: dev.blues.io",
        "• Blues Notehub: notehub.io",
        "",
        "Support:",
        "• GitHub Issues",
        "• Blues Wireless Forums"
    ])
    
    # Save presentation
    output_path = SCREENSHOT_BASE / OUTPUT_FILE
    prs.save(str(output_path))
    print(f"Presentation created successfully: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    
    return str(output_path)

if __name__ == "__main__":
    create_presentation()
